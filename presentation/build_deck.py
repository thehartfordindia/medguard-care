"""
Generate a polished, client-ready MedGuard AI presentation as MedGuard_AI.pptx.

Run:
    python presentation/build_deck.py

Output:
    presentation/MedGuard_AI.pptx

Hartford brand palette
----------------------
Hartford Red    #E4002B
Deep Navy       #0A2540
Slate Ink       #1F2937
Cool Gray       #6B7280
Mist            #F3F4F6
Accent Teal     #14B8A6
Risk Amber      #F59E0B
Risk Crimson    #DC2626
Soft White      #FFFFFF
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ─── Brand palette ────────────────────────────────────────
HARTFORD_RED  = RGBColor(0xE4, 0x00, 0x2B)
NAVY          = RGBColor(0x0A, 0x25, 0x40)
INK           = RGBColor(0x1F, 0x29, 0x37)
GRAY          = RGBColor(0x6B, 0x72, 0x80)
MIST          = RGBColor(0xF3, 0xF4, 0xF6)
TEAL          = RGBColor(0x14, 0xB8, 0xA6)
AMBER         = RGBColor(0xF5, 0x9E, 0x0B)
CRIMSON       = RGBColor(0xDC, 0x26, 0x26)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
EMERALD       = RGBColor(0x10, 0xB9, 0x81)


# ─── Slide geometry (16:9) ────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Helpers ──────────────────────────────────────────────
def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, *,
             font="Calibri", size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=16, color=INK, bullet_color=HARTFORD_RED, line_spacing=1.35):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Red square bullet
        bullet = p.add_run()
        bullet.text = "▌  "
        bullet.font.name = "Calibri"
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = bullet_color
        bullet.font.bold = True
        # Body
        body = p.add_run()
        body.text = item
        body.font.name = "Calibri"
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return tb


def add_footer(slide, page_number=None, total=None):
    # Bottom red accent bar
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), NAVY)
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.35), Inches(0.6), Inches(0.35), HARTFORD_RED)
    add_text(slide, Inches(0.7), SLIDE_H - Inches(0.34), Inches(8), Inches(0.32),
             "MedGuard AI  ·  The Hartford  ·  Confidential — for internal demo",
             size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if page_number is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.34),
                 Inches(1.5), Inches(0.32),
                 f"{page_number} / {total}",
                 size=10, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_header(slide, eyebrow, title):
    # Top navy band
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), Inches(1.15), HARTFORD_RED)
    add_text(slide, Inches(0.55), Inches(0.18), Inches(11), Inches(0.32),
             eyebrow.upper(),
             size=11, bold=True, color=RGBColor(0xFF, 0xC1, 0xC9))
    add_text(slide, Inches(0.55), Inches(0.48), Inches(12), Inches(0.6),
             title,
             size=26, bold=True, color=WHITE)


def add_kpi_card(slide, left, top, width, height, value, label, sub, accent=TEAL):
    add_rect(slide, left, top, width, height, WHITE, line=MIST)
    # Top accent stripe
    add_rect(slide, left, top, width, Inches(0.08), accent)
    add_text(slide, left + Inches(0.2), top + Inches(0.25),
             width - Inches(0.4), Inches(0.7),
             value, size=32, bold=True, color=NAVY)
    add_text(slide, left + Inches(0.2), top + Inches(1.05),
             width - Inches(0.4), Inches(0.4),
             label, size=12, bold=True, color=INK)
    add_text(slide, left + Inches(0.2), top + Inches(1.4),
             width - Inches(0.4), Inches(0.5),
             sub, size=10, color=GRAY)


# ═════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═════════════════════════════════════════════════════════════

def slide_cover(prs, total):
    s = add_blank_slide(prs)
    # Full-bleed navy background
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    # Diagonal red accent
    add_rect(s, Inches(0), Inches(0), Inches(0.4), SLIDE_H, HARTFORD_RED)

    # Hartford-style mark
    add_text(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4),
             "THE HARTFORD  ·  AI / ML / DATA SCIENCE",
             size=12, bold=True, color=RGBColor(0xFF, 0xC1, 0xC9))

    # Big title
    add_text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.4),
             "MedGuard AI",
             size=72, bold=True, color=WHITE)

    add_text(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(0.7),
             "Intelligent Claim Intelligence — powered by AI-Augmented Risk Mitigation",
             size=22, color=RGBColor(0xCB, 0xD5, 0xE1))

    # Divider line
    add_rect(s, Inches(0.9), Inches(4.0), Inches(2.4), Inches(0.05), HARTFORD_RED)

    # Value-prop strip
    add_text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.5),
             "A reusable AI pattern for Middle & Large Commercial and Group Benefits",
             size=16, bold=True, color=WHITE)

    # Author block
    add_text(s, Inches(0.9), Inches(5.4), Inches(11), Inches(0.4),
             "Presented by  ·  [Your Name]",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(5.85), Inches(11), Inches(0.4),
             "Hackathon 2026  ·  Internal Demo",
             size=12, color=GRAY)

    add_footer(s, 1, total)


def slide_agenda(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "Agenda", "What we'll cover in the next five minutes")

    items = [
        ("01", "Our line of business",      "Why Middle & Large Commercial cares about claim intelligence"),
        ("02", "The opportunity",            "The dollar value of catching risk earlier"),
        ("03", "MedGuard AI architecture",   "Four modules, one pipeline, zero vendor lock-in"),
        ("04", "Who uses it & how",          "Personas, day-in-the-life, and three concrete use cases"),
        ("05", "Live demonstration",         "JARVIS · GARMA · NotebookLM · multilingual voice"),
        ("06", "Hartford applications & path to production", "From hackathon prototype to enterprise deployment"),
    ]
    top = Inches(1.6)
    for i, (num, title, sub) in enumerate(items):
        y = top + Inches(0.78 * i)
        # Number block
        add_rect(s, Inches(0.7), y, Inches(0.85), Inches(0.65), HARTFORD_RED)
        add_text(s, Inches(0.7), y, Inches(0.85), Inches(0.65),
                 num, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title + sub
        add_text(s, Inches(1.8), y + Inches(0.02), Inches(11), Inches(0.35),
                 title, size=18, bold=True, color=NAVY)
        add_text(s, Inches(1.8), y + Inches(0.4), Inches(11), Inches(0.3),
                 sub, size=12, color=GRAY)

    add_footer(s, page, total)


def slide_business(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "Our line of business",
               "Middle & Large Commercial and Group Benefits at The Hartford")

    # Three columns
    col_y = Inches(1.6)
    col_w = Inches(4.0)
    col_h = Inches(4.4)
    gap = Inches(0.25)
    start_x = Inches(0.55)

    cols = [
        ("Middle & Large Commercial",
         "Property · Liability · Workers' Comp",
         [
             "Mid-market employers ($25M – $1B revenue)",
             "Fortune-1000 risk programs",
             "Loss-control & claim-cost reduction",
             "Medical evidence drives WC severity",
         ],
         HARTFORD_RED),
        ("Group Benefits",
         "Disability · Absence · Group Life",
         [
             "#2 U.S. carrier in group disability",
             "Short-term, long-term, FMLA, leave",
             "Return-to-work outcomes",
             "Clinical evidence reviewed by nurses",
         ],
         NAVY),
        ("Specialty & Life",
         "Underwriting · Specialty risk",
         [
             "APS / MIB / IME document workflows",
             "Mortality and morbidity scoring",
             "Manual review = bottleneck",
             "Opportunity for AI acceleration",
         ],
         TEAL),
    ]
    for i, (title, sub, bullets, accent) in enumerate(cols):
        x = start_x + i * (col_w + gap)
        add_rect(s, x, col_y, col_w, col_h, WHITE, line=MIST)
        add_rect(s, x, col_y, col_w, Inches(0.12), accent)
        add_text(s, x + Inches(0.25), col_y + Inches(0.3),
                 col_w - Inches(0.5), Inches(0.5),
                 title, size=18, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), col_y + Inches(0.85),
                 col_w - Inches(0.5), Inches(0.4),
                 sub, size=12, color=GRAY)
        add_bullets(s, x + Inches(0.25), col_y + Inches(1.4),
                    col_w - Inches(0.5), Inches(2.8),
                    bullets, size=13, line_spacing=1.45)

    # Closing strip
    add_rect(s, Inches(0.55), Inches(6.25), SLIDE_W - Inches(1.1), Inches(0.6), MIST)
    add_text(s, Inches(0.75), Inches(6.32), SLIDE_W - Inches(1.5), Inches(0.5),
             "All three lines share one truth — the signal is buried in unstructured medical text.",
             size=14, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


def slide_problem(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "The opportunity", "What earlier risk detection is worth")

    # Big number block
    add_rect(s, Inches(0.55), Inches(1.6), Inches(6.0), Inches(4.3), NAVY)
    add_text(s, Inches(0.7), Inches(1.85), Inches(5.7), Inches(0.5),
             "INDUSTRY BENCHMARK", size=11, bold=True, color=RGBColor(0xFF, 0xC1, 0xC9))
    add_text(s, Inches(0.7), Inches(2.4), Inches(5.7), Inches(1.5),
             "$15 – 30M",
             size=80, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(4.1), Inches(5.7), Inches(0.7),
             "Annual value of a 1% reduction in",
             size=18, color=RGBColor(0xCB, 0xD5, 0xE1))
    add_text(s, Inches(0.7), Inches(4.5), Inches(5.7), Inches(0.7),
             "short-term to long-term disability conversion",
             size=18, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(5.3), Inches(5.7), Inches(0.5),
             "Source: industry analysis, top-3 U.S. group disability carrier",
             size=10, color=GRAY)

    # Right column — pain points
    add_text(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.6),
             "Today's reality", size=22, bold=True, color=NAVY)
    pains = [
        "Nurse case managers read records sequentially",
        "Drug-interaction red flags surface after escalation",
        "Underwriters spend 30–60 min per APS letter",
        "Multilingual members get English-only service",
        "Risk indicators sit in unstructured text",
    ]
    add_bullets(s, Inches(7.0), Inches(2.4), Inches(5.8), Inches(3.0),
                pains, size=15, line_spacing=1.5)

    # Bottom callout
    add_rect(s, Inches(7.0), Inches(5.55), Inches(5.8), Inches(0.95), HARTFORD_RED)
    add_text(s, Inches(7.2), Inches(5.65), Inches(5.5), Inches(0.85),
             "MedGuard AI gives every reviewer an AI co-pilot —\nrisk surfaces in seconds, not days.",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


def slide_architecture(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "Architecture",
               "Four modules · one pipeline · zero proprietary APIs")

    # Pipeline visual — five stages
    stages = [
        ("📄  Data",        "FHIR-style records",        TEAL),
        ("🤖  JARVIS",      "NL understanding",          NAVY),
        ("📊  NotebookLM",  "Vector search / RAG",       HARTFORD_RED),
        ("🛡️  GARMA",      "Risk scoring engine",       AMBER),
        ("📈  Cockpit",     "Streamlit dashboard",       EMERALD),
    ]
    box_w = Inches(2.25)
    box_h = Inches(1.4)
    gap = Inches(0.15)
    total_w = len(stages) * box_w + (len(stages) - 1) * gap
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(1.7)

    for i, (title, sub, color) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, y, box_w, box_h, color)
        add_text(s, x + Inches(0.1), y + Inches(0.2),
                 box_w - Inches(0.2), Inches(0.55),
                 title, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), y + Inches(0.85),
                 box_w - Inches(0.2), Inches(0.4),
                 sub, size=12, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Arrow between
        if i < len(stages) - 1:
            arr_x = x + box_w + Inches(0.0)
            add_text(s, arr_x, y + Inches(0.5),
                     gap, Inches(0.4),
                     "›", size=24, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Module table below
    add_text(s, Inches(0.55), Inches(3.55), Inches(12), Inches(0.5),
             "What each module does",
             size=18, bold=True, color=NAVY)

    rows = [
        ("📄  Data Engine",   "FHIR / HL7 inspired", "250 synthetic patients · Pydantic-validated"),
        ("🤖  JARVIS",       "AI-Augmented Engineering", "Plain-English queries · clinical taxonomy match"),
        ("🛡️  GARMA",       "GenAI Risk Mitigation", "6-factor risk score · drug interactions · anomalies"),
        ("📊  NotebookLM",   "Google NotebookLM-style", "FAISS + sentence-transformers · semantic search"),
        ("📈  Cockpit",      "Streamlit",          "KPIs · alerts · 19-language voice · clinician reports"),
    ]
    row_y = Inches(4.05)
    row_h = Inches(0.42)
    for i, (mod, inspired, role) in enumerate(rows):
        y_r = row_y + i * row_h
        bg = MIST if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.55), y_r, SLIDE_W - Inches(1.1), row_h, bg)
        add_text(s, Inches(0.7), y_r + Inches(0.06), Inches(2.6), Inches(0.35),
                 mod, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.4), y_r + Inches(0.06), Inches(3.0), Inches(0.35),
                 inspired, size=11, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(6.5), y_r + Inches(0.06), Inches(6.3), Inches(0.35),
                 role, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


def slide_demo_overview(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "Live demonstration  ·  Part 1",
               "Population cockpit — every claim, every risk, one screen")

    # KPI cards mimicking the dashboard
    kpis = [
        ("250",  "Patients monitored",     "synthetic FHIR records",     TEAL),
        ("47",   "Critical alerts",        "GARMA-flagged today",         CRIMSON),
        ("18",   "Drug interactions",      "severe / moderate severity",  AMBER),
        ("32 %", "Average risk score",     "across full population",      NAVY),
    ]
    card_w = Inches(2.95)
    card_h = Inches(1.95)
    start_x = Inches(0.55)
    y = Inches(1.6)
    for i, (val, label, sub, accent) in enumerate(kpis):
        x = start_x + i * (card_w + Inches(0.15))
        add_kpi_card(s, x, y, card_w, card_h, val, label, sub, accent)

    # Demo callouts
    add_text(s, Inches(0.55), Inches(3.85), Inches(12), Inches(0.5),
             "What you'll see on screen", size=18, bold=True, color=NAVY)
    callouts = [
        "Glassmorphism KPI cards — at-a-glance population health",
        "Live critical-alerts feed — SOC-style risk console",
        "Risk-distribution donut + age vs risk scatter",
        "One-click drill-down to any flagged patient",
    ]
    add_bullets(s, Inches(0.55), Inches(4.4), Inches(12), Inches(2.5),
                callouts, size=15, line_spacing=1.45)

    add_footer(s, page, total)


def slide_demo_jarvis(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "Live demonstration  ·  Part 2",
               "JARVIS — natural-language clinical querying with multilingual voice")

    # Left: query example
    add_rect(s, Inches(0.55), Inches(1.6), Inches(6.2), Inches(4.5), MIST)
    add_text(s, Inches(0.75), Inches(1.75), Inches(6), Inches(0.4),
             "USER QUERY", size=11, bold=True, color=GRAY)
    add_text(s, Inches(0.75), Inches(2.1), Inches(6), Inches(0.7),
             "“Find high-risk elderly patients on warfarin.”",
             size=18, bold=True, color=NAVY)

    add_rect(s, Inches(0.75), Inches(2.95), Inches(5.8), Inches(0.04), HARTFORD_RED)

    add_text(s, Inches(0.75), Inches(3.05), Inches(6), Inches(0.4),
             "JARVIS RESPONSE", size=11, bold=True, color=GRAY)
    sample = (
        "✓ Found 10 matching patients\n"
        "✓ 7 flagged HIGH risk by GARMA\n"
        "✓ 3 with active drug-interaction warnings\n"
        "✓ Avg relevance score: 0.87\n\n"
        "Top match — Brian Yang, 71M\n"
        "Insulin · Amiodarone · Warfarin\n"
        "Risk: 91 % · Bleeding-risk interaction"
    )
    add_text(s, Inches(0.75), Inches(3.45), Inches(6), Inches(2.5),
             sample, size=13, color=INK, line_spacing=1.4)

    # Right: capabilities
    add_text(s, Inches(7.1), Inches(1.7), Inches(5.7), Inches(0.5),
             "Capabilities", size=18, bold=True, color=NAVY)
    caps = [
        "Plain-English clinical queries",
        "Keyword + clinical-taxonomy match",
        "Sub-second response · runs offline",
        "Voice replies in 19 languages",
        "Zero LLM API key required for demo",
    ]
    add_bullets(s, Inches(7.1), Inches(2.25), Inches(5.7), Inches(2.6),
                caps, size=14, line_spacing=1.5)

    # Voice strip
    add_rect(s, Inches(7.1), Inches(4.95), Inches(5.7), Inches(1.15), TEAL)
    add_text(s, Inches(7.3), Inches(5.05), Inches(5.4), Inches(0.4),
             "🔊  MULTILINGUAL VOICE", size=11, bold=True, color=WHITE)
    add_text(s, Inches(7.3), Inches(5.4), Inches(5.4), Inches(0.6),
             "Spanish · Mandarin · Hindi · Tagalog · Arabic",
             size=14, bold=True, color=WHITE)
    add_text(s, Inches(7.3), Inches(5.78), Inches(5.4), Inches(0.4),
             "Production: Azure Speech Service (HIPAA-compliant)",
             size=11, color=WHITE)

    add_footer(s, page, total)


def slide_demo_garma(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "Live demonstration  ·  Part 3",
               "GARMA — explainable risk scoring you can defend in front of a clinician")

    # Six-factor model
    factors = [
        ("Age",                    "demographic baseline"),
        ("Polypharmacy",           "≥ 5 active medications"),
        ("Abnormal labs",          "out-of-range biomarkers"),
        ("Drug interactions",      "severity-weighted pairs"),
        ("Chronic conditions",     "comorbidity load"),
        ("Vital-sign trend",       "recent deterioration"),
    ]
    add_text(s, Inches(0.55), Inches(1.55), Inches(12), Inches(0.5),
             "Six clinical factors, transparently weighted",
             size=18, bold=True, color=NAVY)

    cell_w = Inches(4.0)
    cell_h = Inches(0.95)
    cell_y = Inches(2.15)
    cell_x = Inches(0.55)
    for i, (name, desc) in enumerate(factors):
        col = i % 3
        row = i // 3
        x = cell_x + col * (cell_w + Inches(0.15))
        y = cell_y + row * (cell_h + Inches(0.18))
        add_rect(s, x, y, cell_w, cell_h, WHITE, line=MIST)
        add_rect(s, x, y, Inches(0.12), cell_h, HARTFORD_RED)
        add_text(s, x + Inches(0.25), y + Inches(0.12),
                 cell_w - Inches(0.4), Inches(0.4),
                 name, size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), y + Inches(0.5),
                 cell_w - Inches(0.4), Inches(0.4),
                 desc, size=11, color=GRAY)

    # Bottom: drug-interaction example
    add_rect(s, Inches(0.55), Inches(4.4), SLIDE_W - Inches(1.1), Inches(2.45), NAVY)
    add_text(s, Inches(0.75), Inches(4.5), Inches(12), Inches(0.5),
             "GARMA flagged · severity HIGH",
             size=12, bold=True, color=RGBColor(0xFF, 0xC1, 0xC9))
    add_text(s, Inches(0.75), Inches(4.85), Inches(12), Inches(0.6),
             "Warfarin  +  Aspirin",
             size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.75), Inches(5.4), Inches(12), Inches(0.5),
             "Risk: severe bleeding · GI hemorrhage",
             size=14, color=RGBColor(0xCB, 0xD5, 0xE1))

    # Recommendation block
    add_rect(s, Inches(7.0), Inches(5.95), Inches(5.8), Inches(0.78), HARTFORD_RED)
    add_text(s, Inches(7.2), Inches(6.0), Inches(5.5), Inches(0.7),
             "RECOMMENDATION → Review anticoagulation;\nconsider PPI; confirm INR within 7 days.",
             size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


# ═════════════════════════════════════════════════════════════
# NEW: PERSONAS — who actually uses this
# ═════════════════════════════════════════════════════════════
def slide_personas(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "Who uses MedGuard AI",
               "Four personas across Hartford — every one of them reads medical text today")

    personas = [
        ("🩺", "Nurse Case Manager",   "Group Benefits — LTD",
         "Today: reads 40-page records sequentially.",
         "With MedGuard: ranked work queue + auto-flagged drug interactions.",
         HARTFORD_RED),
        ("💼", "Claim Handler",         "Workers' Compensation",
         "Today: no clinical training, must set reserves.",
         "With MedGuard: plain-English risk score + 'why' explanation.",
         NAVY),
        ("📋", "Underwriter",           "Life / Specialty",
         "Today: 30 – 60 min per APS letter, manual.",
         "With MedGuard: ask APS questions, get cited paragraph answers.",
         TEAL),
        ("📞", "Member-Services Agent", "Group Benefits / WC",
         "Today: English-only callbacks, escalate to interpreter.",
         "With MedGuard: instant voice reply in 19 languages.",
         AMBER),
    ]

    card_w = Inches(6.2)
    card_h = Inches(2.45)
    gap = Inches(0.18)
    start_x = Inches(0.55)
    start_y = Inches(1.55)

    for i, (icon, role, line, today, future, accent) in enumerate(personas):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        add_rect(s, x, y, card_w, card_h, WHITE, line=MIST)
        add_rect(s, x, y, Inches(0.18), card_h, accent)

        # Icon + role
        add_text(s, x + Inches(0.35), y + Inches(0.18),
                 Inches(0.7), Inches(0.6),
                 icon, size=28, bold=True, color=accent)
        add_text(s, x + Inches(1.15), y + Inches(0.2),
                 card_w - Inches(1.4), Inches(0.4),
                 role, size=16, bold=True, color=NAVY)
        add_text(s, x + Inches(1.15), y + Inches(0.6),
                 card_w - Inches(1.4), Inches(0.35),
                 line, size=11, color=GRAY)

        # Today vs future
        add_text(s, x + Inches(0.35), y + Inches(1.1),
                 Inches(1.0), Inches(0.3),
                 "TODAY", size=9, bold=True, color=GRAY)
        add_text(s, x + Inches(1.4), y + Inches(1.08),
                 card_w - Inches(1.7), Inches(0.4),
                 today, size=11, color=INK)

        add_text(s, x + Inches(0.35), y + Inches(1.7),
                 Inches(1.0), Inches(0.3),
                 "WITH AI", size=9, bold=True, color=accent)
        add_text(s, x + Inches(1.4), y + Inches(1.68),
                 card_w - Inches(1.7), Inches(0.5),
                 future, size=11, bold=True, color=NAVY)

    add_footer(s, page, total)


# ═════════════════════════════════════════════════════════════
# NEW: DAY-IN-THE-LIFE — Priya, the LTD nurse case manager
# ═════════════════════════════════════════════════════════════
def slide_day_in_the_life(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "A day in the life",
               "Priya, an LTD nurse case manager — Monday morning, 8:00 a.m.")

    steps = [
        ("8:00", "Opens the cockpit",
         "Sees 250 claimants in her book, 47 critical alerts overnight, 18 drug interactions caught.",
         "5 seconds (vs 45 minutes in Excel today)."),
        ("8:01", "Asks JARVIS in plain English",
         "“Find high-risk elderly patients on warfarin.”  →  10 ranked patients with risk scores.",
         "No SQL, no query builder, no analyst needed."),
        ("8:02", "Drills into the worst patient (GARMA)",
         "Brian Yang, 71M, risk 91 %.  GARMA radar shows polypharmacy + drug interaction as drivers.",
         "Decision-grade explanation, not a black-box score."),
        ("8:03", "Reads the recommendation",
         "Warfarin + Aspirin → severe bleeding risk.  GARMA recommends INR check + PPI within 7 days.",
         "She picks up the phone and calls the treating physician."),
        ("8:05", "Documents and moves on",
         "Downloads the clinician report, attaches it to the claim file, audit-trail recorded.",
         "Next claim — total elapsed time: 5 minutes."),
    ]

    top = Inches(1.55)
    for i, (time, action, detail, value) in enumerate(steps):
        y = top + Inches(1.05 * i)

        # Time chip
        add_rect(s, Inches(0.55), y, Inches(1.0), Inches(0.92), HARTFORD_RED)
        add_text(s, Inches(0.55), y, Inches(1.0), Inches(0.92),
                 time, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Body card
        add_rect(s, Inches(1.65), y, SLIDE_W - Inches(2.2), Inches(0.92), MIST)
        add_text(s, Inches(1.85), y + Inches(0.06),
                 Inches(11), Inches(0.32),
                 action, size=13, bold=True, color=NAVY)
        add_text(s, Inches(1.85), y + Inches(0.36),
                 Inches(11), Inches(0.3),
                 detail, size=10, color=INK)
        add_text(s, Inches(1.85), y + Inches(0.62),
                 Inches(11), Inches(0.28),
                 "→  " + value, size=10, bold=True, color=HARTFORD_RED, font="Calibri")

    add_footer(s, page, total)


# ═════════════════════════════════════════════════════════════
# NEW: CUSTOMER USE CASES — three concrete examples across LOBs
# ═════════════════════════════════════════════════════════════
def slide_use_cases(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "How customers use it",
               "Three concrete examples — same product, three lines of business")

    examples = [
        ("USE CASE 1",
         "Group Benefits  ·  STD-to-LTD triage",
         "An LTD claim costs ~$80,000 average. We need to spot the convertors on day 1.",
         "Nurse types:  “Show STD claims week-3 with opioids and no return-to-work plan.”",
         "JARVIS returns 14 claimants → GARMA scores conversion risk → top 3 get nurse intervention.",
         "1 % conversion reduction = $15 – 30 M / year.",
         HARTFORD_RED),
        ("USE CASE 2",
         "Workers' Compensation  ·  catastrophic-claim watch",
         "Opioid + benzodiazepine combinations are the #1 driver of WC claim escalation.",
         "Claim handler asks:  “Which open WC claims have opioid + benzo prescriptions?”",
         "GARMA flags 9 claims with HIGH severity drug interactions → escalate to medical director.",
         "Earlier intervention → ~25 % shorter claim duration on flagged cases.",
         NAVY),
        ("USE CASE 3",
         "Life / Specialty Underwriting  ·  APS acceleration",
         "Underwriters lose 30 – 60 minutes per Attending Physician Statement.",
         "Underwriter asks NotebookLM:  “Any history of cardiac events in this APS?”",
         "RAG returns the exact paragraph from the document with citation → bind decision in 5 min.",
         "Throughput +25 % on APS-driven cases without adding headcount.",
         TEAL),
    ]

    top = Inches(1.55)
    box_h = Inches(1.7)
    for i, (label, title, problem, action, outcome, value, accent) in enumerate(examples):
        y = top + i * (box_h + Inches(0.1))

        # Left accent panel
        add_rect(s, Inches(0.55), y, Inches(1.7), box_h, accent)
        add_text(s, Inches(0.6), y + Inches(0.2),
                 Inches(1.6), Inches(0.4),
                 label, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(0.6), y + Inches(0.7),
                 Inches(1.6), Inches(0.9),
                 value, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Right content
        add_rect(s, Inches(2.35), y, SLIDE_W - Inches(2.9), box_h, WHITE, line=MIST)
        add_text(s, Inches(2.55), y + Inches(0.1),
                 Inches(10.5), Inches(0.4),
                 title, size=15, bold=True, color=NAVY)

        add_text(s, Inches(2.55), y + Inches(0.55),
                 Inches(0.7), Inches(0.3),
                 "Pain", size=9, bold=True, color=GRAY)
        add_text(s, Inches(3.3), y + Inches(0.53),
                 Inches(9.7), Inches(0.35),
                 problem, size=10, color=INK)

        add_text(s, Inches(2.55), y + Inches(0.92),
                 Inches(0.7), Inches(0.3),
                 "Action", size=9, bold=True, color=GRAY)
        add_text(s, Inches(3.3), y + Inches(0.9),
                 Inches(9.7), Inches(0.35),
                 action, size=10, color=INK)

        add_text(s, Inches(2.55), y + Inches(1.3),
                 Inches(0.7), Inches(0.3),
                 "Result", size=9, bold=True, color=accent)
        add_text(s, Inches(3.3), y + Inches(1.28),
                 Inches(9.7), Inches(0.35),
                 outcome, size=10, bold=True, color=NAVY)

    add_footer(s, page, total)


# ═════════════════════════════════════════════════════════════
# (existing) HARTFORD MAPPING — kept after the use-case slides
# ═════════════════════════════════════════════════════════════
def slide_hartford_mapping(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "Hartford applications",
               "Same architecture · three immediate Hartford use cases")

    # Header row
    headers = ["MedGuard module", "Group Benefits LTD", "Workers' Comp", "Life / Specialty UW"]
    col_xs = [Inches(0.55), Inches(3.6), Inches(7.0), Inches(10.3)]
    col_ws = [Inches(2.95), Inches(3.3), Inches(3.2), Inches(2.45)]
    header_y = Inches(1.55)
    add_rect(s, Inches(0.55), header_y, SLIDE_W - Inches(1.1), Inches(0.5), NAVY)
    for x, w, h in zip(col_xs, col_ws, headers):
        add_text(s, x, header_y, w, Inches(0.5),
                 h, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("🤖  JARVIS",
         "“Show LTD claimants with no RTW plan”",
         "“Find catastrophic claims on opioids”",
         "Underwriter NL search over APS"),
        ("🛡️  GARMA",
         "STD → LTD conversion risk score",
         "Opioid / polypharmacy red flags",
         "Mortality composite score"),
        ("📊  NotebookLM",
         "RAG over IME and APS reports",
         "Treatment-note semantic search",
         "MIB / APS PDF retrieval"),
        ("📈  Cockpit",
         "Nurse case-manager triage queue",
         "Claim-handler severity dashboard",
         "Underwriter decision support"),
    ]
    row_y = Inches(2.1)
    row_h = Inches(0.7)
    for i, row in enumerate(rows):
        y = row_y + i * row_h
        bg = MIST if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.55), y, SLIDE_W - Inches(1.1), row_h, bg)
        for x, w, val in zip(col_xs, col_ws, row):
            bold = (x == col_xs[0])
            color = NAVY if bold else INK
            add_text(s, x, y + Inches(0.1), w, row_h - Inches(0.2),
                     val, size=12, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE)

    # Value strip
    add_text(s, Inches(0.55), Inches(5.25), Inches(12), Inches(0.5),
             "Estimated value (industry benchmarks)",
             size=18, bold=True, color=NAVY)
    vals = [
        ("$15 – 30M / yr", "from a 1% LTD-duration reduction",         HARTFORD_RED),
        ("25 % faster",     "underwriter case review cycle time",       NAVY),
        ("5 – 10 %",        "call-center deflection via voice agent",   TEAL),
    ]
    box_w = Inches(4.0)
    bx = Inches(0.55)
    by = Inches(5.85)
    for i, (big, small, color) in enumerate(vals):
        x = bx + i * (box_w + Inches(0.15))
        add_rect(s, x, by, box_w, Inches(1.0), color)
        add_text(s, x + Inches(0.2), by + Inches(0.1),
                 box_w - Inches(0.4), Inches(0.45),
                 big, size=20, bold=True, color=WHITE)
        add_text(s, x + Inches(0.2), by + Inches(0.55),
                 box_w - Inches(0.4), Inches(0.4),
                 small, size=11, color=WHITE)

    add_footer(s, page, total)


def slide_path_to_prod(prs, total, page):
    s = add_blank_slide(prs)
    add_header(s, "Path to production",
               "From hackathon prototype to Hartford-grade deployment")

    items = [
        ("01", "Data sources",
         "Replace synthetic generator with Guidewire ClaimCenter, APS feeds, Rx data"),
        ("02", "Enterprise AI services",
         "Azure OpenAI + Azure Translator / Speech (existing Hartford licenses)"),
        ("03", "Security & compliance",
         "Hartford SSO · PHI guardrails · audit logging · HITRUST · HIPAA"),
        ("04", "Hosting",
         "Azure App Service / AKS behind Zscaler · zero data egress"),
        ("05", "Model governance",
         "Register GARMA in Hartford AI Model Risk inventory · NAIC-aligned monitoring"),
    ]
    top = Inches(1.7)
    for i, (num, title, sub) in enumerate(items):
        y = top + Inches(0.95 * i)
        add_rect(s, Inches(0.55), y, Inches(0.85), Inches(0.78), NAVY)
        add_text(s, Inches(0.55), y, Inches(0.85), Inches(0.78),
                 num, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(1.5), y, SLIDE_W - Inches(2.05), Inches(0.78), MIST)
        add_text(s, Inches(1.7), y + Inches(0.08), Inches(11), Inches(0.35),
                 title, size=15, bold=True, color=NAVY)
        add_text(s, Inches(1.7), y + Inches(0.42), Inches(11), Inches(0.35),
                 sub, size=12, color=INK)

    add_footer(s, page, total)


def slide_close(prs, total, page):
    s = add_blank_slide(prs)
    # Background
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.4), SLIDE_H, HARTFORD_RED)

    add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.5),
             "MEDGUARD AI", size=14, bold=True, color=RGBColor(0xFF, 0xC1, 0xC9))
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.4),
             "A pattern, not just a prototype.",
             size=48, bold=True, color=WHITE)

    add_rect(s, Inches(0.9), Inches(2.95), Inches(2.4), Inches(0.05), HARTFORD_RED)

    points = [
        ("✓", "Four modules · one architecture · zero vendor lock-in"),
        ("✓", "End-to-end working demo in five minutes"),
        ("✓", "Re-applicable across Disability, Workers' Comp, Underwriting"),
        ("✓", "Clear path to Azure-hosted, HIPAA-compliant production"),
    ]
    for i, (mark, text) in enumerate(points):
        y = Inches(3.3) + Inches(0.55 * i)
        add_text(s, Inches(0.9), y, Inches(0.45), Inches(0.4),
                 mark, size=20, bold=True, color=TEAL)
        add_text(s, Inches(1.4), y, Inches(11.5), Inches(0.4),
                 text, size=18, color=WHITE)

    add_text(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.5),
             "Thank you.", size=24, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
             "Questions?", size=18, color=RGBColor(0xCB, 0xD5, 0xE1))

    add_footer(s, page, total)


# ═════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,            # 1
        slide_agenda,           # 2
        slide_business,         # 3
        slide_problem,          # 4
        slide_architecture,     # 5
        slide_personas,         # 6  — NEW: who uses it
        slide_demo_overview,    # 7
        slide_demo_jarvis,      # 8
        slide_demo_garma,       # 9
        slide_day_in_the_life,  # 10 — NEW: walkthrough
        slide_use_cases,        # 11 — NEW: 3 concrete customer use cases
        slide_hartford_mapping, # 12
        slide_path_to_prod,     # 13
        slide_close,            # 14
    ]
    total = len(builders)

    # Cover takes (prs, total) — others take (prs, total, page)
    slide_cover(prs, total)
    for i, fn in enumerate(builders[1:], start=2):
        fn(prs, total, i)

    out = Path(__file__).parent / "MedGuard_AI.pptx"
    prs.save(str(out))
    print(f"✓ Deck written to {out}")


if __name__ == "__main__":
    build()
